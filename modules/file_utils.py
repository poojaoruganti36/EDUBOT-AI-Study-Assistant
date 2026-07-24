import json
from io import BytesIO

import streamlit as st
from docx import Document
from pptx import Presentation

import config
from modules.image_utils import extract_text_from_image
from modules.pdf_utils import extract_pdf_text
from modules.rag import create_index, retrieve, split_text


def _read_text_file(uploaded_file) -> str:
    return uploaded_file.getvalue().decode("utf-8", errors="ignore")


def _read_csv_file(uploaded_file) -> str:
    content = uploaded_file.getvalue().decode("utf-8", errors="ignore")
    return content.replace(",", " | ")


def _read_json_file(uploaded_file) -> str:
    parsed = json.loads(uploaded_file.getvalue().decode("utf-8", errors="ignore"))
    return json.dumps(parsed, indent=2, ensure_ascii=False)


def _read_docx_file(uploaded_file) -> str:
    document = Document(BytesIO(uploaded_file.getvalue()))
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    return "\n".join(paragraphs)


def _read_pptx_file(uploaded_file) -> str:
    presentation = Presentation(BytesIO(uploaded_file.getvalue()))
    slides_text: list[str] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        shape_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                cleaned = shape.text.strip()
                if cleaned:
                    shape_text.append(cleaned)
        if shape_text:
            slides_text.append(f"Slide {slide_number}\n" + "\n".join(shape_text))

    return "\n\n".join(slides_text)


def extract_text_from_uploads(uploaded_files) -> tuple[str, list[str]]:
    if not uploaded_files:
        return "", []

    sections: list[str] = []
    sources: list[str] = []

    for uploaded_file in uploaded_files:
        suffix = uploaded_file.name.lower().split(".")[-1]
        try:
            if suffix == "pdf":
                text = extract_pdf_text([uploaded_file])
            elif suffix in {"txt", "md"}:
                text = _read_text_file(uploaded_file)
            elif suffix == "csv":
                text = _read_csv_file(uploaded_file)
            elif suffix == "json":
                text = _read_json_file(uploaded_file)
            elif suffix == "docx":
                text = _read_docx_file(uploaded_file)
            elif suffix == "pptx":
                text = _read_pptx_file(uploaded_file)
            elif suffix in {"png", "jpg", "jpeg"}:
                text = extract_text_from_image(uploaded_file)
            else:
                text = ""

            if text.strip():
                sections.append(f"Source: {uploaded_file.name}\n{text.strip()}")
                sources.append(uploaded_file.name)
        except Exception as exc:
            st.warning(f"Could not read {uploaded_file.name}: {exc}")

    return "\n\n".join(sections), sources


def build_knowledge_base(text: str):
    chunks = split_text(text)
    index = create_index(chunks)

    def retriever(question: str) -> str:
        results = retrieve(question, index, chunks)
        joined = "\n\n".join(results)
        return joined[: config.MAX_CONTEXT_CHARS]

    return {"chunks": chunks, "index": index, "retriever": retriever}
